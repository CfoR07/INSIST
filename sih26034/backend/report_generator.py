import os
from typing import Dict, Any, List
from datetime import datetime
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors

def generate_html_report(inspection: Dict[str, Any], images: List[Dict[str, Any]], facts: List[Dict[str, Any]], results: List[Dict[str, Any]], reviews: List[Dict[str, Any]]) -> str:
    reviews_by_res = {r["compliance_result_id"]: r for r in reviews}
    pass_cnt = sum(1 for r in results if r["status"] == "PASS")
    fail_cnt = sum(1 for r in results if r["status"] == "FAIL")
    unc_cnt = sum(1 for r in results if r["status"] in ["UNCERTAIN", "CONFLICT"])
    na_cnt = sum(1 for r in results if r["status"] == "NOT_APPLICABLE")
    
    rows_html = ""
    for r in results:
        rev = reviews_by_res.get(r["id"])
        officer_text = f"<span class='text-emerald-700 font-bold'>Confirmed: {rev['decision']}</span>" if rev else "<span class='text-gray-400 italic'>None</span>"
        status_badge = {
            "PASS": "<span class='px-2.5 py-1 bg-emerald-100 text-emerald-800 rounded font-bold text-xs'>PASS</span>",
            "FAIL": "<span class='px-2.5 py-1 bg-rose-100 text-rose-800 rounded font-bold text-xs'>FAIL</span>",
            "UNCERTAIN": "<span class='px-2.5 py-1 bg-amber-100 text-amber-800 rounded font-bold text-xs'>UNCERTAIN</span>",
            "CONFLICT": "<span class='px-2.5 py-1 bg-purple-100 text-purple-800 rounded font-bold text-xs'>CONFLICT</span>",
            "NOT_APPLICABLE": "<span class='px-2.5 py-1 bg-gray-100 text-gray-800 rounded font-bold text-xs'>N/A</span>"
        }.get(r["status"], r["status"])
        
        rows_html += f"""
        <tr class="border-b border-gray-200 hover:bg-gray-50">
            <td class="py-3 px-4 font-mono text-xs font-bold text-gray-600">{r.get('rule_id')}</td>
            <td class="py-3 px-4 font-medium text-gray-900 text-sm">{r.get('requirement')}<br><span class="text-xs text-gray-400">{r.get('source_reference')}</span></td>
            <td class="py-3 px-4 font-mono text-xs text-gray-800">{r.get('observed_value') or '—'}</td>
            <td class="py-3 px-4">{status_badge}</td>
            <td class="py-3 px-4 text-xs text-gray-600">{r.get('reason')}</td>
            <td class="py-3 px-4 text-xs">{officer_text}</td>
        </tr>
        """

    return f"""<!DOCTYPE html>
<html lang="en">
<head>
    <meta charset="UTF-8">
    <title>Inspection Report - {inspection.get('id')}</title>
    <script src="https://cdn.tailwindcss.com"></script>
</head>
<body class="bg-gray-100 p-8 font-sans antialiased text-gray-800">
    <div class="max-w-5xl mx-auto bg-white shadow-xl rounded-xl border border-gray-200 overflow-hidden print:shadow-none print:border-none">
        <div class="bg-slate-900 text-white p-6 flex justify-between items-center">
            <div>
                <div class="flex items-center gap-2">
                    <span class="bg-blue-600 text-xs px-2 py-0.5 rounded font-bold uppercase tracking-wider">Official Inspection Audit</span>
                    <span class="text-slate-400 text-xs font-mono">{datetime.now().strftime('%d-%b-%Y %H:%M:%S')}</span>
                </div>
                <h1 class="text-2xl font-bold mt-1">Legal Metrology Packaged Commodities Report</h1>
                <p class="text-sm text-slate-300">Department of Consumer Affairs • Legal Metrology Enforcement Division</p>
            </div>
            <div class="text-right">
                <div class="text-xs text-slate-400 uppercase tracking-wider">Inspection ID</div>
                <div class="text-2xl font-mono font-bold text-blue-400">{inspection.get('id')}</div>
            </div>
        </div>

        <div class="p-6 bg-slate-50 border-b border-gray-200 grid grid-cols-2 sm:grid-cols-4 gap-4 text-sm">
            <div>
                <div class="text-xs text-gray-500 font-medium uppercase">Product Name</div>
                <div class="font-bold text-gray-800">{inspection.get('product_name')}</div>
            </div>
            <div>
                <div class="text-xs text-gray-500 font-medium uppercase">Category</div>
                <div class="font-bold text-gray-800">{inspection.get('category')}</div>
            </div>
            <div>
                <div class="text-xs text-gray-500 font-medium uppercase">Enforcement Officer</div>
                <div class="font-bold text-gray-800">{inspection.get('officer_id')}</div>
            </div>
            <div>
                <div class="text-xs text-gray-500 font-medium uppercase">Inspection Status</div>
                <div class="font-bold text-blue-700">{inspection.get('status')}</div>
            </div>
        </div>

        <div class="p-6 border-b border-gray-200 grid grid-cols-4 gap-4">
            <div class="bg-emerald-50 border border-emerald-200 p-4 rounded-lg text-center">
                <div class="text-2xl font-black text-emerald-700">{pass_cnt}</div>
                <div class="text-xs font-bold text-emerald-800 uppercase mt-1">Rules Passed</div>
            </div>
            <div class="bg-rose-50 border border-rose-200 p-4 rounded-lg text-center">
                <div class="text-2xl font-black text-rose-700">{fail_cnt}</div>
                <div class="text-xs font-bold text-rose-800 uppercase mt-1">Violations Found</div>
            </div>
            <div class="bg-amber-50 border border-amber-200 p-4 rounded-lg text-center">
                <div class="text-2xl font-black text-amber-700">{unc_cnt}</div>
                <div class="text-xs font-bold text-amber-800 uppercase mt-1">Review Cases</div>
            </div>
            <div class="bg-gray-50 border border-gray-200 p-4 rounded-lg text-center">
                <div class="text-2xl font-black text-gray-700">{na_cnt}</div>
                <div class="text-xs font-bold text-gray-800 uppercase mt-1">Exempt / NA</div>
            </div>
        </div>

        <div class="p-6">
            <h2 class="text-lg font-bold text-gray-900 mb-4 flex items-center justify-between">
                <span>Statutory Declarations Audit Breakdown</span>
                <span class="text-xs font-normal text-gray-500">Deterministic Rule Engine v2024</span>
            </h2>
            <div class="overflow-x-auto">
                <table class="w-full text-left border-collapse">
                    <thead>
                        <tr class="bg-slate-100 text-slate-700 uppercase font-semibold text-xs border-b border-gray-300">
                            <th class="py-3 px-4">Rule ID</th>
                            <th class="py-3 px-4">Statutory Requirement</th>
                            <th class="py-3 px-4">Observed Value</th>
                            <th class="py-3 px-4">Verdict</th>
                            <th class="py-3 px-4">Engine Validation Reason</th>
                            <th class="py-3 px-4">Officer Action</th>
                        </tr>
                    </thead>
                    <tbody>
                        {rows_html}
                    </tbody>
                </table>
            </div>
        </div>

        <div class="p-6 bg-gray-50 border-t border-gray-200 flex justify-between items-center text-xs text-gray-500">
            <div>
                Generated automatically by SIH26034 Pre-Packed Commodity Inspection Engine.<br>
                Core principle: <em>"AI extracts observations. Deterministic code enforces compliance."</em>
            </div>
            <div class="text-right border-t-2 border-gray-400 pt-2 w-48">
                <div class="font-bold text-gray-800">{inspection.get('officer_id')}</div>
                <div>Authorized Officer Signature</div>
            </div>
        </div>
    </div>
</body>
</html>"""

def generate_pdf_report(inspection: Dict[str, Any], results: List[Dict[str, Any]], output_path: str):
    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    doc = SimpleDocTemplate(output_path, pagesize=letter, leftMargin=36, rightMargin=36, topMargin=36, bottomMargin=36)
    styles = getSampleStyleSheet()
    story = []
    
    title_style = ParagraphStyle(name='TitleStyle', parent=styles['Heading1'], fontSize=16, leading=20, textColor=colors.HexColor('#0f172a'))
    story.append(Paragraph(f"<b>SIH26034 Legal Metrology Inspection Report</b>", title_style))
    story.append(Paragraph(f"Inspection ID: {inspection.get('id')} | Product: {inspection.get('product_name')} ({inspection.get('category')})", styles['Normal']))
    story.append(Paragraph(f"Officer ID: {inspection.get('officer_id')} | Date: {datetime.now().strftime('%Y-%m-%d %H:%M')}", styles['Normal']))
    story.append(Spacer(1, 14))
    
    data = [["Rule ID", "Statutory Requirement", "Observed Fact", "Verdict", "Validation Details"]]
    for r in results:
        data.append([
            r.get("rule_id", ""),
            r.get("requirement", "")[:35],
            str(r.get("observed_value", ""))[:20],
            r.get("status", ""),
            r.get("reason", "")[:50]
        ])
        
    t = Table(data, colWidths=[65, 140, 95, 65, 175])
    t.setStyle(TableStyle([
        ('BACKGROUND', (0,0), (-1,0), colors.HexColor('#1e293b')),
        ('TEXTCOLOR', (0,0), (-1,0), colors.whitesmoke),
        ('ALIGN', (0,0), (-1,-1), 'LEFT'),
        ('FONTNAME', (0,0), (-1,0), 'Helvetica-Bold'),
        ('FONTSIZE', (0,0), (-1,0), 9),
        ('BOTTOMPADDING', (0,0), (-1,0), 6),
        ('BACKGROUND', (0,1), (-1,-1), colors.HexColor('#f8fafc')),
        ('GRID', (0,0), (-1,-1), 0.5, colors.HexColor('#cbd5e1')),
        ('FONTSIZE', (0,1), (-1,-1), 8),
    ]))
    story.append(t)
    story.append(Spacer(1, 20))
    story.append(Paragraph("<i>Compliance decision is determined strictly by deterministic Python rules operating on structured extracted facts.</i>", styles['Italic']))
    
    doc.build(story)
    return output_path

def generate_pptx_report(inspection: Dict[str, Any], results: List[Dict[str, Any]], output_path: str):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN

    os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # 1. Slide 1: Title Slide (Liquid Dark Noir)
    blank_slide_layout = prs.slide_layouts[6]
    slide1 = prs.slides.add_slide(blank_slide_layout)
    
    # Background shape
    bg = slide1.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bg.line.fill.background()

    txBox = slide1.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.3), Inches(3.0))
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "INSIST"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)

    p2 = tf.add_paragraph()
    p2.text = "Legal Metrology Packaged Commodities Inspection Audit"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(96, 165, 250)
    p2.space_before = Pt(12)

    p3 = tf.add_paragraph()
    p3.text = f"Inspection ID: {inspection.get('id')}  |  Product: {inspection.get('product_name', 'Commodity')}  |  Officer: {inspection.get('officer_id', 'OFFICER-001')}"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(203, 213, 225)
    p3.space_before = Pt(20)

    # 2. Slide 2: Compliance Results Matrix
    slide2 = prs.slides.add_slide(blank_slide_layout)
    bg2 = slide2.shapes.add_shape(1, 0, 0, prs.slide_width, prs.slide_height)
    bg2.fill.solid()
    bg2.fill.fore_color.rgb = RGBColor(15, 23, 42)
    bg2.line.fill.background()

    # Header
    tx_header = slide2.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.7), Inches(0.8))
    tf_h = tx_header.text_frame
    p_h = tf_h.paragraphs[0]
    p_h.text = "LMPC Rules 2011 Compliance Audit Results"
    p_h.font.size = Pt(22)
    p_h.font.bold = True
    p_h.font.color.rgb = RGBColor(255, 255, 255)

    # Table
    rows = len(results) + 1
    cols = 5
    table_shape = slide2.shapes.add_table(rows, cols, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    table = table_shape.table

    headers = ["Rule ID", "Statutory Requirement", "Observed Fact", "Verdict", "Validation Reason"]
    widths = [Inches(1.5), Inches(3.2), Inches(2.0), Inches(1.2), Inches(3.8)]
    for i, w in enumerate(widths):
        table.columns[i].width = w

    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(30, 41, 59)
        for p in cell.text_frame.paragraphs:
            p.font.bold = True
            p.font.size = Pt(11)
            p.font.color.rgb = RGBColor(255, 255, 255)

    for row_idx, r in enumerate(results, start=1):
        vals = [
            r.get("rule_id", ""),
            r.get("requirement", "")[:45],
            str(r.get("observed_value", ""))[:25],
            r.get("status", ""),
            r.get("reason", "")[:60]
        ]
        for col_idx, val in enumerate(vals):
            cell = table.cell(row_idx, col_idx)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(24, 32, 47) if row_idx % 2 == 0 else RGBColor(30, 41, 59)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9.5)
                if col_idx == 3:
                    if val == "PASS":
                        p.font.color.rgb = RGBColor(52, 211, 153)
                        p.font.bold = True
                    elif val == "FAIL":
                        p.font.color.rgb = RGBColor(248, 113, 113)
                        p.font.bold = True
                    else:
                        p.font.color.rgb = RGBColor(251, 191, 36)
                else:
                    p.font.color.rgb = RGBColor(241, 245, 249)

    prs.save(output_path)
    return output_path

